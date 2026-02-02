"""Document loading and preprocessing for various file formats."""

import os
from pathlib import Path
from typing import List, Dict, Any, Optional
import logging

import pypdf
from pptx import Presentation
from bs4 import BeautifulSoup
import html2text

from app.core.config import settings

logger = logging.getLogger(__name__)


class DocumentLoader:
    """Loads and preprocesses documents from various formats."""
    
    def __init__(self):
        self.html_converter = html2text.HTML2Text()
        self.html_converter.ignore_links = False
        self.html_converter.ignore_images = True
        self.html_converter.body_width = 0  # Don't wrap lines
    
    def load_document(
        self,
        file_path: Optional[str] = None,
        file_url: Optional[str] = None,
        content: Optional[str] = None,
        file_type: str = "text"
    ) -> Dict[str, Any]:
        """
        Load a document from file path, URL, or raw content.
        
        Returns:
            Dictionary with 'content' (str) and 'metadata' (dict)
        """
        if file_path:
            return self._load_from_file(file_path, file_type)
        elif file_url:
            return self._load_from_url(file_url, file_type)
        elif content:
            return {
                "content": self._clean_text(content),
                "metadata": {"source": "direct_input", "type": file_type}
            }
        else:
            raise ValueError("Must provide file_path, file_url, or content")
    
    def _load_from_file(self, file_path: str, file_type: str) -> Dict[str, Any]:
        """Load document from local file path."""
        path = Path(file_path)
        
        if not path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")
        
        if file_type == "pdf":
            return self._load_pdf(file_path)
        elif file_type == "pptx":
            return self._load_pptx(file_path)
        elif file_type == "html":
            return self._load_html(file_path)
        elif file_type == "text":
            return self._load_text(file_path)
        else:
            raise ValueError(f"Unsupported file type: {file_type}")
    
    def _load_pdf(self, file_path: str) -> Dict[str, Any]:
        """Load and extract text from PDF."""
        try:
            content_parts = []
            metadata = {
                "source": file_path,
                "type": "pdf",
                "pages": []
            }
            
            with open(file_path, "rb") as file:
                pdf_reader = pypdf.PdfReader(file)
                metadata["total_pages"] = len(pdf_reader.pages)
                
                for page_num, page in enumerate(pdf_reader.pages, start=1):
                    text = page.extract_text()
                    if text.strip():
                        content_parts.append(text)
                        metadata["pages"].append(page_num)
            
            content = "\n\n".join(content_parts)
            content = self._clean_text(content)
            
            logger.info(f"Loaded PDF: {file_path} ({len(pdf_reader.pages)} pages)")
            
            return {
                "content": content,
                "metadata": metadata
            }
        except Exception as e:
            logger.error(f"Error loading PDF {file_path}: {str(e)}")
            raise

    def _load_pptx(self, file_path: str) -> Dict[str, Any]:
        """Load and extract text from PowerPoint (.pptx) file."""
        try:
            content_parts = []
            metadata = {"source": file_path, "type": "pptx", "slides": []}
            prs = Presentation(file_path)
            metadata["total_slides"] = len(prs.slides)

            for slide_num, slide in enumerate(prs.slides, start=1):
                slide_texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_texts.append(shape.text.strip())
                if slide_texts:
                    slide_content = "\n".join(slide_texts)
                    content_parts.append(f"Slide {slide_num}:\n{slide_content}")
                    metadata["slides"].append(slide_num)

            content = "\n\n".join(content_parts)
            content = self._clean_text(content)

            logger.info(f"Loaded PPTX: {file_path} ({len(prs.slides)} slides)")
            return {"content": content, "metadata": metadata}
        except Exception as e:
            logger.error(f"Error loading PPTX {file_path}: {str(e)}")
            raise
    
    def _load_html(self, file_path: str) -> Dict[str, Any]:
        """Load and extract text from HTML file."""
        try:
            with open(file_path, "r", encoding="utf-8") as file:
                html_content = file.read()
            
            soup = BeautifulSoup(html_content, "lxml")
            
            # Remove navigation, headers, footers, and boilerplate
            for element in soup.find_all(["nav", "header", "footer", "script", "style"]):
                element.decompose()
            
            # Extract main content
            main_content = soup.find("main") or soup.find("article") or soup.find("body")
            
            if main_content:
                text = self.html_converter.handle(str(main_content))
            else:
                text = self.html_converter.handle(html_content)
            
            # Extract title if available
            title = soup.find("title")
            title_text = title.get_text() if title else None
            
            content = self._clean_text(text)
            
            metadata = {
                "source": file_path,
                "type": "html",
                "title": title_text
            }
            
            logger.info(f"Loaded HTML: {file_path}")
            
            return {
                "content": content,
                "metadata": metadata
            }
        except Exception as e:
            logger.error(f"Error loading HTML {file_path}: {str(e)}")
            raise
    
    def _load_text(self, file_path: str) -> Dict[str, Any]:
        """Load plain text file."""
        try:
            with open(file_path, "r", encoding="utf-8") as file:
                content = file.read()
            
            content = self._clean_text(content)
            
            metadata = {
                "source": file_path,
                "type": "text"
            }
            
            logger.info(f"Loaded text file: {file_path}")
            
            return {
                "content": content,
                "metadata": metadata
            }
        except Exception as e:
            logger.error(f"Error loading text file {file_path}: {str(e)}")
            raise
    
    def _load_from_url(self, url: str, file_type: str) -> Dict[str, Any]:
        """Load document from URL (placeholder for future implementation)."""
        # TODO: Implement URL fetching with requests/httpx
        raise NotImplementedError("URL loading not yet implemented")
    
    def _clean_text(self, text: str) -> str:
        """
        Clean text by removing excessive whitespace, references, etc.
        """
        # Remove excessive whitespace
        lines = text.split("\n")
        cleaned_lines = []
        
        for line in lines:
            line = line.strip()
            # Skip empty lines and very short lines that are likely navigation
            if len(line) < 3:
                continue
            # Skip lines that look like references (e.g., "[1]", "(Smith et al., 2020)")
            if line.startswith("[") and line.endswith("]"):
                continue
            # Skip lines that are mostly numbers/punctuation
            if sum(c.isalnum() for c in line) < len(line) * 0.3:
                continue
            cleaned_lines.append(line)
        
        # Join with single newlines
        cleaned = "\n".join(cleaned_lines)
        
        # Remove multiple consecutive newlines
        while "\n\n\n" in cleaned:
            cleaned = cleaned.replace("\n\n\n", "\n\n")
        
        return cleaned.strip()

